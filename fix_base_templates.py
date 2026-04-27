"""One-shot: split the legacy single base template into Indian + Western
variants, fixing the cross-vertical leftover phrases in each.

Inputs : _base_template.pptx (legacy, hand-edited)
Outputs: _base_template_indian.pptx, _base_template_western.pptx

Slide 5 point 2 and slide 8 point 1 in the legacy base ship with stray
electronics + grocery examples ("16GM RAM", "I want a laptop good for
Photoshop", "plan my meals") that do not belong in a fashion default. We
replace them with vertical-coherent phrases — Indian on one base, Western
on the other.

Run once:
    python3 fix_base_templates.py
"""

from pathlib import Path
import shutil

from pptx import Presentation

from pptx_helpers import replace_in_shape


BASE_DIR = Path(__file__).parent
LEGACY = BASE_DIR / "_base_template.pptx"
INDIAN = BASE_DIR / "_base_template_indian.pptx"
WESTERN = BASE_DIR / "_base_template_western.pptx"


# Smart quotes match the master template's existing characters.
LDQUO, RDQUO, RSQUO = "“", "”", "’"

INDIAN_REPLACEMENTS = {
    # Slide 5 (idx 4) — point 2 "Conversational search needs a knowledge
    # graph to map intent to keywords"
    4: [
        (f"{LDQUO}16GM RAM{RDQUO}", f"{LDQUO}cotton kurta{RDQUO}"),
        (
            f"{LDQUO}I want a laptop good for Photoshop{RDQUO}",
            f"{LDQUO}I want a saree gift for my mother-in-law{RSQUO}s "
            f"birthday{RDQUO}",
        ),
    ],
    # Slide 8 (idx 7) — point 1 "AI Sales Associate"
    7: [
        (
            f"{LDQUO}plan my meals{RDQUO}",
            f"{LDQUO}what to wear to a sangeet{RDQUO}",
        ),
    ],
}

WESTERN_REPLACEMENTS = {
    4: [
        (f"{LDQUO}16GM RAM{RDQUO}", f"{LDQUO}summer floral dress{RDQUO}"),
        (
            f"{LDQUO}I want a laptop good for Photoshop{RDQUO}",
            f"{LDQUO}I want an outfit for my college reunion{RDQUO}",
        ),
    ],
    7: [
        (
            f"{LDQUO}plan my meals{RDQUO}",
            f"{LDQUO}what to wear to a friend{RSQUO}s wedding{RDQUO}",
        ),
    ],
}


def apply_replacements(pptx_path, replacements_per_slide):
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    for slide_idx, swaps in replacements_per_slide.items():
        for shape in slides[slide_idx].shapes:
            replace_in_shape(shape, swaps)
    prs.save(str(pptx_path))


def main():
    if not LEGACY.exists():
        raise SystemExit(
            f"Expected legacy base at {LEGACY}. This is a one-shot script — "
            f"if the split has already been done, inspect _base_template_*.pptx "
            f"directly."
        )

    shutil.copy(LEGACY, INDIAN)
    apply_replacements(INDIAN, INDIAN_REPLACEMENTS)
    print(f"Wrote {INDIAN.name}")

    shutil.copy(LEGACY, WESTERN)
    apply_replacements(WESTERN, WESTERN_REPLACEMENTS)
    print(f"Wrote {WESTERN.name}")

    LEGACY.unlink()
    print(f"Removed legacy {LEGACY.name}")


if __name__ == "__main__":
    main()
