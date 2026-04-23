"""Build category-specific generic pitch deck templates.

Usage:
    python build_category_templates.py --all
    python build_category_templates.py --category electronics
    python build_category_templates.py --category bags,shoes,jewellery
"""

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from pptx_helpers import (
    A,
    find_shape_by_text,
    replace_bullet_column,
    replace_in_shape,
    set_shape_text,
)
from category_content import CATEGORIES
from graph_generator import generate_graph


BASE_DIR = Path(__file__).parent
TEMPLATE = BASE_DIR / "_base_template.pptx"
OUTPUT_DIR = BASE_DIR / "output"
GRAPHS_DIR = BASE_DIR / "graphs"


def build_template(category_key):
    cat = CATEGORIES[category_key]
    label = cat["label"]
    out_path = OUTPUT_DIR / f"SS_Generic_{label}.pptx"

    graph_path = GRAPHS_DIR / f"{category_key}_graph.png"
    if not graph_path.exists():
        generate_graph(cat["graph_definition"], graph_path)

    shutil.copy(TEMPLATE, out_path)
    prs = Presentation(str(out_path))
    slides = list(prs.slides)

    # --- Slide 3: competitors -----------------------------------------------
    if cat["competitors_global"] or cat["competitors_india"]:
        s3 = slides[2]
        global_col = None
        india_col = None
        for shape in s3.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if "Amazon Rufus" in txt and "Walmart Sparky" in txt:
                global_col = shape
            elif "Myntra" in txt or "Aza/Pernia" in txt or "Meesho" in txt:
                india_col = shape
        if global_col and cat["competitors_global"]:
            replace_bullet_column(global_col, cat["competitors_global"])
        if india_col and cat["competitors_india"]:
            replace_bullet_column(india_col, cat["competitors_india"])

    # --- Cross-slide text swaps (slides 4, 5, 8, 12) -----------------------
    for slide_idx, swaps in cat.get("text_swaps", {}).items():
        if swaps:
            for shape in slides[slide_idx].shapes:
                replace_in_shape(shape, swaps)

    # --- Slide 10: upcoming offerings ---------------------------------------
    if cat.get("slide10_offerings"):
        s10 = slides[9]
        for cfg in cat["slide10_offerings"]:
            title_shape = None
            desc_shape = None
            for shape in s10.shapes:
                if not shape.has_text_frame:
                    continue
                txt = shape.text_frame.text.strip()
                if not txt:
                    continue
                if title_shape is None and txt == cfg["title_match"]:
                    title_shape = shape
                elif desc_shape is None and cfg["desc_match"] in txt:
                    desc_shape = shape
            if title_shape:
                set_shape_text(title_shape, cfg["title"])
            if desc_shape:
                set_shape_text(desc_shape, cfg["desc"])

    # --- Slide 11: replace knowledge graph image ----------------------------
    s11 = slides[10]
    for shape in s11.shapes:
        if shape.shape_type == 13:  # PICTURE
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            s11.shapes.add_picture(
                str(graph_path), left, top, width=width, height=height,
            )
            break

    # --- Slide 13: comparison table example queries -------------------------
    if cat.get("comparison_examples"):
        s13 = slides[12]
        cmp_tbl = None
        for shape in s13.shapes:
            if shape.has_table:
                cmp_tbl = shape.table
                break
        if cmp_tbl is not None:
            for r_idx, row in enumerate(cmp_tbl.rows):
                if r_idx in cat["comparison_examples"]:
                    cells = list(row.cells)
                    cell = cells[0]
                    paras = cell.text_frame._txBody.findall(f"{A}p")
                    if len(paras) >= 2:
                        query_para = paras[1]
                        runs = query_para.findall(f"{A}r")
                        if runs:
                            first_run = runs[0]
                            first_run.find(f"{A}t").text = (
                                cat["comparison_examples"][r_idx]
                            )
                            for r in runs[1:]:
                                query_para.remove(r)

    prs.save(str(out_path))
    print(f"Built {out_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Build category-specific generic pitch deck templates."
    )
    parser.add_argument(
        "--category", type=str, default=None,
        help="Comma-separated category keys (e.g. electronics,furniture) or 'all'",
    )
    parser.add_argument(
        "--all", action="store_true",
        help="Build all categories",
    )
    args = parser.parse_args()

    OUTPUT_DIR.mkdir(exist_ok=True)
    GRAPHS_DIR.mkdir(exist_ok=True)

    if args.all or args.category == "all":
        targets = list(CATEGORIES.keys())
    elif args.category:
        targets = [c.strip() for c in args.category.split(",")]
        for t in targets:
            if t not in CATEGORIES:
                parser.error(f"Unknown category: {t}. Choose from {list(CATEGORIES.keys())}")
    else:
        parser.error("Specify --all or --category <name>")

    for cat_key in targets:
        build_template(cat_key)

    print(f"\nDone. Built {len(targets)} template(s) in {OUTPUT_DIR}/")


if __name__ == "__main__":
    main()
